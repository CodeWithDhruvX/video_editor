import os
import json
import tkinter as tk
from tkinter import filedialog, ttk, messagebox, scrolledtext
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import markdown2
import re

# Mobile-optimized themes with better contrast
THEMES = {
    "Light": {"bg_color": "#FFFFFF", "text_color": "#000000", "accent": "#007ACC"},
    "Dark": {"bg_color": "#1E1E1E", "text_color": "#FFFFFF", "accent": "#4FC3F7"},
    "Blue": {"bg_color": "#0D47A1", "text_color": "#FFFFFF", "accent": "#81D4FA"},
    "Green": {"bg_color": "#1B5E20", "text_color": "#FFFFFF", "accent": "#81C784"},
    "Purple": {"bg_color": "#4A148C", "text_color": "#FFFFFF", "accent": "#CE93D8"}
}

# Programming language extensions mapping
LANGUAGE_EXTENSIONS = {
    'python': '.py', 'javascript': '.js', 'java': '.java', 'cpp': '.cpp', 'c++': '.cpp',
    'c': '.c', 'csharp': '.cs', 'c#': '.cs', 'php': '.php', 'ruby': '.rb', 'go': '.go',
    'rust': '.rs', 'swift': '.swift', 'kotlin': '.kt', 'scala': '.scala', 'r': '.r',
    'matlab': '.m', 'sql': '.sql', 'html': '.html', 'css': '.css', 'typescript': '.ts',
    'bash': '.sh', 'shell': '.sh', 'powershell': '.ps1', 'dockerfile': 'Dockerfile',
    'yaml': '.yml', 'yml': '.yml', 'json': '.json', 'xml': '.xml', 'markdown': '.md',
    'tex': '.tex', 'latex': '.tex', 'vim': '.vim', 'lua': '.lua', 'perl': '.pl',
    'dart': '.dart', 'haskell': '.hs', 'clojure': '.clj', 'elixir': '.ex', 'erlang': '.erl',
    'f#': '.fs', 'groovy': '.groovy', 'julia': '.jl', 'objc': '.m', 'pascal': '.pas',
    'fortran': '.f90', 'assembly': '.asm', 'vhdl': '.vhd', 'verilog': '.v', 'txt': '.txt'
}


class FinalPPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("📱 Final Mobile PPT Generator - Multi JSON Support + Code Files")
        self.root.geometry("700x850") # Increased height for better visibility
        self.all_topics_data = []
        self.loaded_files = []
        
        # Setup UI
        self.setup_ui()
        
    def setup_ui(self):
        # Main frame with scrollbar
        main_canvas = tk.Canvas(self.root)
        scrollbar = ttk.Scrollbar(self.root, orient="vertical", command=main_canvas.yview)
        scrollable_frame = ttk.Frame(main_canvas)

        scrollable_frame.bind(
            "<Configure>",
            lambda e: main_canvas.configure(scrollregion=main_canvas.bbox("all"))
        )

        main_canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        main_canvas.configure(yscrollcommand=scrollbar.set)
        
        main_frame = ttk.Frame(scrollable_frame, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # Title
        title_label = ttk.Label(main_frame, text="📱 Final Mobile PPT Generator", 
                                font=("Arial", 18, "bold"))
        title_label.pack(pady=(0, 5))
        
        subtitle_label = ttk.Label(main_frame, text="Multi-JSON Support -  Organized Folders -  Mobile Optimized -  Auto Code Files", 
                                   font=("Arial", 10), foreground="gray")
        subtitle_label.pack(pady=(0, 20))
        
        # Feature highlight
        feature_frame = ttk.Frame(main_frame)
        feature_frame.pack(fill=tk.X, pady=(0, 15))
        
        feature_label = ttk.Label(feature_frame, text="🆕 NEW: Automatically creates separate code files for slides with 'slide_type': 'code'", 
                                  font=("Arial", 9, "bold"), foreground="green")
        feature_label.pack()
        
        # JSON Input Section
        json_frame = ttk.LabelFrame(main_frame, text="📄 JSON Input Management", padding="15")
        json_frame.pack(fill=tk.X, pady=(0, 20))
        
        # Buttons frame
        buttons_frame = ttk.Frame(json_frame)
        buttons_frame.pack(fill=tk.X, pady=(0, 15))
        
        load_multiple_btn = tk.Button(buttons_frame, text="📂 Load Multiple JSON Files", 
                                      font=("Arial", 11, "bold"), bg="#28a745", fg="white",
                                      command=self.load_multiple_json_files, cursor="hand2")
        load_multiple_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        clear_btn = tk.Button(buttons_frame, text="🗑️ Clear All", 
                              font=("Arial", 11), bg="#dc3545", fg="white",
                              command=self.clear_all_files, cursor="hand2")
        clear_btn.pack(side=tk.RIGHT)
        
        # JSON Text Input Area
        text_input_frame = ttk.LabelFrame(json_frame, text="✏️ Paste JSON Content", padding="10")
        text_input_frame.pack(fill=tk.BOTH, expand=True, pady=(10, 0))
        
        instruction_label = ttk.Label(text_input_frame, 
                                      text="Paste your JSON array content below and click 'Add from Text':",
                                      font=("Arial", 10))
        instruction_label.pack(anchor=tk.W, pady=(0, 5))
        
        self.json_text_area = scrolledtext.ScrolledText(text_input_frame, 
                                                        height=8, 
                                                        width=60,
                                                        font=("Consolas", 10),
                                                        wrap=tk.WORD)
        self.json_text_area.pack(fill=tk.BOTH, expand=True, pady=(0, 10))
        
        placeholder_text = '''[
 {
   "topic": "Go Basics",
   "slides": [
     {
       "title": "Hello World in Go",
       "slide_type": "code",
       "language": "go",
       "content": "package main\\n\\nimport \\"fmt\\"\\n\\nfunc main() {\\n    fmt.Println(\\"Hello, World!\\")\\n}"
     },
     {
       "title": "What are Variables?",
       "content": "Variables are containers for storing data values.\\n\\n* `var` declares 1 or more variables.\\n* `:=` for short declaration.\\n* **Strongly typed!**",
       "slide_type": "text"
     }
   ]
 }
]'''
        self.json_text_area.insert(tk.END, placeholder_text)
        self.json_text_area.bind("<FocusIn>", self.clear_placeholder)
        
        text_buttons_frame = ttk.Frame(text_input_frame)
        text_buttons_frame.pack(fill=tk.X)
        
        add_text_btn = tk.Button(text_buttons_frame, text="➕ Add from Text", 
                                 font=("Arial", 11, "bold"), bg="#17a2b8", fg="white",
                                 command=self.add_from_text, cursor="hand2")
        add_text_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        clear_text_btn = tk.Button(text_buttons_frame, text="🧹 Clear Text", 
                                   font=("Arial", 10), bg="#6c757d", fg="white",
                                   command=self.clear_text_area, cursor="hand2")
        clear_text_btn.pack(side=tk.LEFT)
        
        # Files list frame
        files_frame = ttk.Frame(json_frame)
        files_frame.pack(fill=tk.BOTH, expand=True, pady=(15, 0))
        
        ttk.Label(files_frame, text="Loaded Sources:", font=("Arial", 10, "bold")).pack(anchor=tk.W, pady=(0, 5))
        
        files_list_frame = ttk.Frame(files_frame)
        files_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.files_listbox = tk.Listbox(files_list_frame, height=4, font=("Arial", 9))
        files_scroll = ttk.Scrollbar(files_list_frame, orient="vertical", command=self.files_listbox.yview)
        self.files_listbox.configure(yscrollcommand=files_scroll.set)
        self.files_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        files_scroll.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Configuration Section
        config_frame = ttk.LabelFrame(main_frame, text="⚙️ Configuration", padding="15")
        config_frame.pack(fill=tk.X, pady=(0, 20))
        
        ttk.Label(config_frame, text="🎨 Select Theme:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        self.theme_var = tk.StringVar(value="Dark")
        self.theme_dropdown = ttk.Combobox(config_frame, textvariable=self.theme_var, 
                                           values=list(THEMES.keys()), state="readonly", width=20)
        self.theme_dropdown.pack(anchor=tk.W, pady=(5, 15))
        
        ttk.Label(config_frame, text="💻 Code File Generation:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        self.generate_code_files_var = tk.BooleanVar(value=True)
        code_check = ttk.Checkbutton(config_frame, text="Generate separate code files for 'slide_type': 'code'", 
                                     variable=self.generate_code_files_var)
        code_check.pack(anchor=tk.W, pady=(5, 15))
        
        ttk.Label(config_frame, text="📁 Starting Folder Number:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        self.folder_number_var = tk.StringVar(value="1")
        folder_number_entry = ttk.Entry(config_frame, textvariable=self.folder_number_var, width=10)
        folder_number_entry.pack(anchor=tk.W, pady=(5, 15))
        
        ttk.Label(config_frame, text="📋 Folder Naming:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        self.naming_var = tk.StringVar(value="numbered")
        naming_radio_frame = ttk.Frame(config_frame)
        naming_radio_frame.pack(anchor=tk.W, pady=(5, 15))
        
        ttk.Radiobutton(naming_radio_frame, text="01. Topic Name", 
                        variable=self.naming_var, value="numbered").pack(side=tk.LEFT, padx=(0, 15))
        ttk.Radiobutton(naming_radio_frame, text="Topic Name Only", 
                        variable=self.naming_var, value="name_only").pack(side=tk.LEFT)
        
        ttk.Label(config_frame, text="📂 Output Directory:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        dir_select_frame = ttk.Frame(config_frame)
        dir_select_frame.pack(fill=tk.X, pady=(5, 10))
        
        self.output_dir_var = tk.StringVar(value=os.getcwd())
        output_dir_entry = ttk.Entry(dir_select_frame, textvariable=self.output_dir_var, state="readonly")
        output_dir_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))
        
        browse_btn = tk.Button(dir_select_frame, text="Browse", 
                               font=("Arial", 9), command=self.browse_output_directory)
        browse_btn.pack(side=tk.RIGHT)
        
        # Topics Preview Section
        preview_frame = ttk.LabelFrame(main_frame, text="📋 All Topics Preview", padding="15")
        preview_frame.pack(fill=tk.X, pady=(0, 20))
        
        self.topics_info_label = ttk.Label(preview_frame, text="No topics loaded", 
                                           font=("Arial", 10), foreground="gray")
        self.topics_info_label.pack(anchor=tk.W, pady=(0,10))
        
        topics_list_frame = ttk.Frame(preview_frame)
        topics_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.topics_listbox = tk.Listbox(topics_list_frame, height=6, font=("Arial", 9))
        topics_scroll = ttk.Scrollbar(topics_list_frame, orient="vertical", command=self.topics_listbox.yview)
        self.topics_listbox.configure(yscrollcommand=topics_scroll.set)
        self.topics_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        topics_scroll.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Generate button
        generate_btn = tk.Button(main_frame, text="🚀 Generate All PPTs with Organized Folders + Code Files", 
                                 font=("Arial", 14, "bold"), bg="#007ACC", fg="white",
                                 command=self.generate_all_presentations, cursor="hand2",
                                 height=2)
        generate_btn.pack(pady=20, fill=tk.X)
        
        self.status_label = tk.Label(main_frame, text="Ready to create organized presentations with automatic code file generation", 
                                     font=("Arial", 10), fg="gray")
        self.status_label.pack(pady=(10, 0))

        main_canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

    def clear_placeholder(self, event):
        current_text = self.json_text_area.get("1.0", tk.END).strip()
        if "Hello World in Go" in current_text:
            self.json_text_area.delete("1.0", tk.END)

    def clear_text_area(self):
        self.json_text_area.delete("1.0", tk.END)

    def add_from_text(self):
        json_text = self.json_text_area.get("1.0", tk.END).strip()
        if not json_text:
            messagebox.showwarning("Warning", "Please paste JSON content in the text area")
            return
        try:
            data = json.loads(json_text)
            if not self.validate_json_structure(data, "Text Input"):
                return
            for topic in data:
                topic["_source_file"] = "Text Input"
            self.all_topics_data.extend(data)
            source_name = f"Text Input ({len(data)} topics)"
            if source_name not in self.loaded_files:
                self.loaded_files.append(source_name)
            self.update_ui_after_load()
            self.json_text_area.delete("1.0", tk.END)
            self.status_label.config(text=f"✅ Added {len(data)} topics from text input", fg="green")
        except json.JSONDecodeError as e:
            messagebox.showerror("JSON Error", f"Invalid JSON format:\n{str(e)}")
            self.status_label.config(text="❌ Invalid JSON format", fg="red")
        except Exception as e:
            messagebox.showerror("Error", f"Error processing text input:\n{str(e)}")
            self.status_label.config(text="❌ Error processing input", fg="red")

    def load_multiple_json_files(self):
        file_paths = filedialog.askopenfilenames(
            filetypes=[("JSON Files", "*.json")], 
            title="Select Multiple JSON Files"
        )
        if not file_paths:
            return

        successful_loads = 0
        failed_loads = []
        for file_path in file_paths:
            try:
                filename = os.path.basename(file_path)
                if filename not in [f.split(' (')[0] for f in self.loaded_files if not f.startswith("Text Input")]:
                    success = self.load_json_data(file_path)
                    if success:
                        successful_loads += 1
                        with open(file_path, "r", encoding="utf-8") as f:
                           file_data = json.load(f)
                        self.loaded_files.append(f"{filename} ({len(file_data)} topics)")
                    else:
                        failed_loads.append(filename)
                else:
                    self.status_label.config(text=f"⚠️ {filename} already loaded", fg="orange")
            except Exception as e:
                failed_loads.append(f"{os.path.basename(file_path)}: {str(e)}")
        
        self.update_ui_after_load()
        if successful_loads > 0:
            msg = f"✅ Loaded {successful_loads} files"
            if failed_loads:
                msg += f", {len(failed_loads)} failed"
            self.status_label.config(text=msg, fg="green" if not failed_loads else "orange")
        elif failed_loads:
            self.status_label.config(text=f"❌ Failed to load {len(failed_loads)} files", fg="red")
            messagebox.showerror("Load Error", f"Failed files:\n" + "\n".join(failed_loads[:3]))

    def load_json_data(self, file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                data = json.load(file)
            if not self.validate_json_structure(data, file_path):
                return False
            for topic in data:
                topic["_source_file"] = os.path.basename(file_path)
            self.all_topics_data.extend(data)
            return True
        except Exception as e:
            messagebox.showerror("JSON Error", f"Error in {os.path.basename(file_path)}:\n{str(e)}")
            return False

    def validate_json_structure(self, data, source):
        source_name = source if isinstance(source, str) and not os.path.exists(source) else os.path.basename(source)
        if not isinstance(data, list):
            messagebox.showerror("Invalid JSON", f"{source_name}: JSON should be a list of topics")
            return False
        for i, topic in enumerate(data):
            if not isinstance(topic, dict):
                messagebox.showerror("Invalid JSON", f"{source_name}: Topic {i+1} should be a dictionary")
                return False
            if "topic" not in topic:
                messagebox.showerror("Invalid JSON", f"{source_name}: Topic {i+1} missing 'topic' field")
                return False
            if "slides" not in topic or not isinstance(topic["slides"], list):
                messagebox.showerror("Invalid JSON", f"{source_name}: Topic {i+1} missing or invalid 'slides' field")
                return False
        return True

    def clear_all_files(self):
        if not self.all_topics_data:
            return
        if messagebox.askyesno("Clear All", "Are you sure you want to clear all loaded content?"):
            self.all_topics_data = []
            self.loaded_files = []
            self.update_ui_after_load()
            self.status_label.config(text="🗑️ All content cleared", fg="orange")

    def update_ui_after_load(self):
        self.files_listbox.delete(0, tk.END)
        for source in self.loaded_files:
            self.files_listbox.insert(tk.END, source)
        
        self.topics_listbox.delete(0, tk.END)
        total_slides = 0
        total_code_slides = 0
        for i, topic in enumerate(self.all_topics_data, 1):
            topic_name = topic.get("topic", f"Topic {i}")
            slides = topic.get("slides", [])
            slides_count = len(slides)
            code_slides = sum(1 for slide in slides if slide.get("slide_type") == "code")
            source_file = topic.get("_source_file", "Unknown")
            total_slides += slides_count
            total_code_slides += code_slides
            display_text = f"{i}. {topic_name} ({slides_count} slides"
            if code_slides > 0:
                display_text += f", {code_slides} code"
            display_text += f") - {source_file}"
            self.topics_listbox.insert(tk.END, display_text)
        
        if self.all_topics_data:
            source_count = len(self.loaded_files)
            info_text = f"📊 Total: {len(self.all_topics_data)} topics, {total_slides} slides"
            if total_code_slides > 0 and self.generate_code_files_var.get():
                info_text += f", {total_code_slides} code files will be generated"
            info_text += f" from {source_count} sources"
            self.topics_info_label.config(text=info_text, foreground="blue")
        else:
            self.topics_info_label.config(text="No topics loaded", foreground="gray")

    def browse_output_directory(self):
        directory = filedialog.askdirectory(title="Select Output Directory")
        if directory:
            self.output_dir_var.set(directory)

    def sanitize_filename(self, filename):
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            filename = filename.replace(char, '_')
        filename = re.sub(r'\s+', ' ', filename.strip())
        return filename[:100]

    def detect_programming_language(self, slide_data):
        if "language" in slide_data and slide_data["language"]:
            lang = slide_data["language"].lower()
            if lang in LANGUAGE_EXTENSIONS:
                return lang

        content = slide_data.get("content", "").strip()
        if not content: return 'txt'
        clean_code = content
        
        if content.startswith("```"):
            lines = content.split('\n')
            first_line = lines[0]
            lang_hint = first_line[3:].strip().lower()
            if lang_hint and lang_hint in LANGUAGE_EXTENSIONS:
                return lang_hint
            if len(lines) > 1 and lines[-1].strip() == "```":
                clean_code = '\n'.join(lines[1:-1])
            else:
                clean_code = '\n'.join(lines[1:])
        
        title = slide_data.get("title", "").lower()
        title_indicators = {
            'python': ['python', 'py', 'django', 'flask'], 'javascript': ['javascript', 'js', 'node', 'react', 'vue', 'es6'],
            'java': ['java', 'spring'], 'cpp': ['c++', 'cpp'], 'c': ['c programming', ' c '],
            'csharp': ['c#', 'csharp', '.net'], 'php': ['php', 'laravel'], 'ruby': ['ruby', 'rails'], 'go': ['golang', 'go'],
            'rust': ['rust'], 'swift': ['swift', 'ios'], 'kotlin': ['kotlin', 'android'], 'sql': ['sql', 'database', 'query'],
            'html': ['html'], 'css': ['css', 'stylesheet', 'style'], 'typescript': ['typescript', 'ts'],
            'bash': ['bash', 'shell', 'script'], 'powershell': ['powershell', 'ps1'], 'yaml': ['yaml', 'yml'],
            'json': ['json'], 'xml': ['xml'], 'dockerfile': ['dockerfile']
        }
        for lang, indicators in title_indicators.items():
            if any(indicator in title for indicator in indicators):
                return lang

        code_lower = clean_code.lower()
        if not code_lower.strip(): return 'txt'

        if 'FROM' in clean_code and 'RUN' in clean_code and ('CMD' in clean_code or 'ENTRYPOINT' in clean_code): return 'dockerfile'
        if '#include' in code_lower and ('<iostream>' in code_lower or 'std::' in code_lower): return 'cpp'
        if '#include' in code_lower and ('<stdio.h>' in code_lower or '<stdlib.h>' in code_lower): return 'c'
        if 'using System;' in clean_code or 'Console.WriteLine' in clean_code or 'namespace' in code_lower: return 'csharp'
        if 'public class' in code_lower and 'public static void main' in code_lower: return 'java'
        if 'package main' in code_lower and 'import "fmt"' in code_lower and 'func main' in code_lower: return 'go'
        if 'fn main' in code_lower and ('println!' in code_lower or 'use std::' in code_lower): return 'rust'
        if re.search(r'\bdef\s+\w+\s*\(.*\):', clean_code) or 'import ' in code_lower: return 'python'
        if '<?php' in code_lower: return 'php'
        if re.search(r'\bfunction\s', clean_code, re.IGNORECASE) or 'console.log' in code_lower or 'const ' in code_lower or 'let ' in code_lower: return 'javascript'
        if 'SELECT' in clean_code.upper() and 'FROM' in clean_code.upper(): return 'sql'
        if '<!DOCTYPE html>' in code_lower or '<html' in code_lower: return 'html'
        if re.search(r'\{[^{}]*:[^{}]*\}', clean_code) and re.search(r'[\w-]+\s*:\s*[\w\s#\'"-]+;', clean_code): return 'css'
        if '#!/bin/bash' in clean_code or '#!/bin/sh' in clean_code or re.search(r'\b(echo|if|fi|then|else)\b', code_lower): return 'bash'
        try:
            if clean_code.strip().startswith(("{", "[")):
                json.loads(clean_code)
                return 'json'
        except (json.JSONDecodeError, TypeError): pass
        return 'txt'

    def get_file_extension(self, language):
        return LANGUAGE_EXTENSIONS.get(language.lower(), '.txt')

    def create_code_file(self, folder_path, topic_name, slide_data, slide_index):
        if not self.generate_code_files_var.get():
            return None
        try:
            language = self.detect_programming_language(slide_data)
            extension = self.get_file_extension(language)
            slide_title = slide_data.get("title", f"slide_{slide_index}")
            sanitized_title = self.sanitize_filename(slide_title)
            
            if language == 'dockerfile':
                filename = "Dockerfile"
            else:
                filename = f"{slide_index:02d}_{sanitized_title}{extension}"
            
            file_path = os.path.join(folder_path, filename)
            
            code_for_file = slide_data.get("content", "")
            stripped_code = code_for_file.strip()
            # FIX #1: Corrected the typo here
            if stripped_code.startswith("```"):
                lines = stripped_code.split('\n')
                start_index = 1
                end_index = len(lines)
                if lines[-1].strip() == "```":
                    end_index -= 1
                code_for_file = '\n'.join(lines[start_index:end_index])
            
            with open(file_path, "w", encoding="utf-8") as f:
                header_title = slide_data.get('title', 'Code Example')
                if language in ['python', 'bash', 'shell', 'ruby', 'perl', 'yaml']: f.write(f"# {header_title}\n# Generated from: {topic_name}\n\n")
                elif language in ['javascript', 'java', 'cpp', 'c', 'csharp', 'php', 'go', 'rust', 'swift', 'kotlin', 'scala', 'css', 'typescript', 'dart', 'groovy']: f.write(f"// {header_title}\n// Generated from: {topic_name}\n\n")
                elif language in ['html', 'xml']: f.write(f"\n\n\n")
                elif language == 'sql': f.write(f"-- {header_title}\n-- Generated from: {topic_name}\n\n")
                f.write(code_for_file)
            return filename
        except Exception as e:
            messagebox.showerror("Error Creating File", f"Could not create code file for slide '{slide_data.get('title', 'Untitled')}':\n{e}")
            return None
            
    def _add_formatted_text(self, shape, text_content, text_color):
        """Adds markdown-formatted text to a shape."""
        # Use a placeholder for newlines that markdown won't strip
        text_content = text_content.replace('\n', '<br>')
        html = markdown2.markdown(text_content, extras=["fenced-code-blocks", "tables", "break-on-newline"])
        soup = BeautifulSoup(html, 'html.parser')

        # Clear any existing text but keep the first paragraph
        shape.text_frame.clear() 
        p = shape.text_frame.paragraphs[0]
        p.text = "" 

        def process_node(node, current_p, is_bold=False):
            if isinstance(node, str):
                run = current_p.add_run()
                run.text = node
                font = run.font
                font.bold = is_bold
                font.color.rgb = text_color
                return current_p

            for child in node.children:
                if isinstance(child, str):
                    run = current_p.add_run()
                    run.text = child
                    font = run.font
                    font.bold = is_bold
                    font.color.rgb = text_color
                elif child.name in ['strong', 'b']:
                    current_p = process_node(child, current_p, is_bold=True)
                elif child.name == 'p':
                    if current_p.text or any(run.text for run in current_p.runs):
                        current_p = shape.text_frame.add_paragraph()
                    current_p = process_node(child, current_p, is_bold=False)
                elif child.name in ['ul', 'ol']:
                    for li in child.find_all('li', recursive=False):
                        if current_p.text or any(run.text for run in current_p.runs):
                           current_p = shape.text_frame.add_paragraph()
                        current_p.text = f"• {li.get_text()}"
                        current_p.level = 1
                        current_p.font.color.rgb = text_color
                elif child.name == 'br':
                    current_p.add_run().text = '\n'
                elif child.name == 'code':
                    run = current_p.add_run()
                    run.text = f" `{child.get_text()}` "
                    font = run.font
                    font.bold = is_bold
                    font.color.rgb = RGBColor(0xE8, 0x3E, 0x82) # Pinkish
                else:
                    current_p = process_node(child, current_p, is_bold)
            return current_p

        process_node(soup, p)


    def generate_all_presentations(self):
        """Generate PPTs for all loaded topics organized in folders with code files"""
        if not self.all_topics_data:
            messagebox.showwarning("Warning", "Please load at least one JSON file or add content from text")
            return
        
        try:
            starting_number = int(self.folder_number_var.get())
        except ValueError:
            messagebox.showerror("Error", "Please enter a valid starting folder number")
            return
        
        output_base_dir = self.output_dir_var.get()
        if not os.path.exists(output_base_dir):
            messagebox.showerror("Error", "Output directory does not exist")
            return
        
        theme_name = self.theme_var.get()
        theme = THEMES.get(theme_name, THEMES["Dark"])
        naming_style = self.naming_var.get()
        
        bg_color = RGBColor.from_string(theme["bg_color"][1:])
        text_color = RGBColor.from_string(theme["text_color"][1:])
        accent_color = RGBColor.from_string(theme["accent"][1:])
        
        try:
            created_folders = []
            total_topics = len(self.all_topics_data)
            total_code_files = 0
            
            for i, topic_data in enumerate(self.all_topics_data):
                self.status_label.config(text=f"⚙️ Generating ({i+1}/{total_topics}): {topic_data.get('topic')}...", fg="blue")
                self.root.update_idletasks()

                folder_number = starting_number + i
                topic_name = topic_data.get("topic", f"Topic {i+1}")
                sanitized_topic = self.sanitize_filename(topic_name)
                
                if naming_style == "numbered":
                    folder_name = f"{folder_number:02d}. {sanitized_topic}"
                else:
                    folder_name = sanitized_topic
                
                folder_path = os.path.join(output_base_dir, folder_name)
                os.makedirs(folder_path, exist_ok=True)
                created_folders.append(folder_name)
                
                slides = topic_data.get("slides", [])
                if self.generate_code_files_var.get():
                    for slide_index, slide_data in enumerate(slides, 1):
                        if slide_data.get("slide_type") == "code":
                            code_filename = self.create_code_file(folder_path, topic_name, slide_data, slide_index)
                            if code_filename:
                                total_code_files += 1
                
                ppt_filename = f"{sanitized_topic}.pptx"
                ppt_path = os.path.join(folder_path, ppt_filename)
                
                prs = Presentation()
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)

                # FIX #2: Set background on the slide master
                slide_master = prs.slide_masters[0]
                fill = slide_master.background.fill
                fill.solid()
                fill.fore_color.rgb = bg_color
                
                blank_slide_layout = prs.slide_layouts[6]

                slide = prs.slides.add_slide(blank_slide_layout)
                title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = topic_name
                p.font.size = Pt(60)
                p.font.bold = True
                p.font.color.rgb = text_color
                p.alignment = PP_ALIGN.CENTER

                for slide_index, slide_data in enumerate(slides, 1):
                    slide = prs.slides.add_slide(blank_slide_layout)

                    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
                    tf_title = title_shape.text_frame
                    p_title = tf_title.paragraphs[0]
                    p_title.text = slide_data.get("title", f"Slide {slide_index}")
                    p_title.font.size = Pt(36)
                    p_title.font.bold = True
                    p_title.font.color.rgb = accent_color
                    p_title.alignment = PP_ALIGN.LEFT
                    
                    content = slide_data.get("content", "")
                    content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(6.5))
                    tf_content = content_shape.text_frame
                    tf_content.word_wrap = True
                    tf_content.auto_size = False

                    if slide_data.get("slide_type") == "code":
                        p_content = tf_content.paragraphs[0]
                        p_content.text = f"Refer to the accompanying code file.\n({self.detect_programming_language(slide_data)})"
                        p_content.font.size = Pt(24)
                        p_content.font.italic = True
                        p_content.font.color.rgb = text_color
                    else:
                        self._add_formatted_text(content_shape, content, text_color)
                        for para in tf_content.paragraphs:
                             if para.font.size is None:
                                  para.font.size = Pt(22)
                
                prs.save(ppt_path)
            
            success_msg = (
                f"🎉 Success! Generated {total_topics} PPTs in {len(created_folders)} folders.\n"
                f"Total code files created: {total_code_files}\n\n"
                f"Output Location: {output_base_dir}"
            )
            messagebox.showinfo("Generation Complete", success_msg)
            self.status_label.config(text=f"✅ Successfully generated {total_topics} presentations!", fg="green")
            
        except Exception as e:
            messagebox.showerror("Generation Error", f"An error occurred during PPT generation:\n{e}")
            self.status_label.config(text=f"❌ Error during generation. Please check console.", fg="red")
            import traceback
            traceback.print_exc()

if __name__ == "__main__":
    root = tk.Tk()
    app = FinalPPTCreatorApp(root)
    root.mainloop()


ppt json:

[
  {
    "topic": "if/elif/else",
    "slides": [
      {
        "title": "Stop Writing Messy if/else Statements in Python!",
        "content": "**Your `if/else` blocks are probably a mess. Here's a simple trick to clean them up!** 🚀\n\nLet's make your code readable and bug-free!",
        "slide_type": "text"
      },
      {
        "title": "Why Use `elif`?",
        "content": "✅ Avoid writing multiple `if` statements that run separately.\n\n✅ `elif` connects conditions for better flow and faster checks.\n\nExample:\n```python\nif score >= 90:\n    print('A grade')\nelif score >= 80:\n    print('B grade')\nelse:\n    print('Try harder!')\n```",
        "slide_type": "code"
      },
      {
        "title": "Common Mistake 🚫",
        "content": "- Using separate `if`s instead of `elif` causes multiple blocks to run.\n- Makes debugging hard and slows down your program.\n\n**Bad:**\n```python\nif score >= 90:\n    print('A grade')\nif score >= 80:\n    print('B grade')\n```",
        "slide_type": "code"
      },
      {
        "title": "Quick Tips ✨",
        "content": "- Use `if` for the first condition.\n- Use `elif` for all middle conditions.\n- Use `else` as a fallback.\n\nThink of it as a flowchart guiding your program.",
        "slide_type": "text"
      },
      {
        "title": "Want Cleaner Code? 👀",
        "content": "Next up: How to use **ternary operators** in Python for one-liners that'll wow your interviewer! Stay tuned! 🔥",
        "slide_type": "text"
      }
    ]
  },
  {
    "topic": "for vs while loops",
    "slides": [
      {
        "title": "Python for vs while Loops: The 60-Second Rule",
        "content": "🤔 Confused when to use `for` or `while`? Here's the fastest way to decide! ⏱️",
        "slide_type": "text"
      },
      {
        "title": "Use `for` When...",
        "content": "- You know **how many times** to repeat.\n- You have a list, range, or collection to loop over.\n\nExample:\n```python\nfor i in range(5):\n    print(i)\n```",
        "slide_type": "code"
      },
      {
        "title": "Use `while` When...",
        "content": "- You want to repeat **until a condition is false**.\n- You don't know how many times it will run.\n\nExample:\n```python\ncount = 0\nwhile count < 5:\n    print(count)\n    count += 1\n```",
        "slide_type": "code"
      },
      {
        "title": "Quick Comparison Table 📊",
        "content": "| Loop Type | When to Use | How it Works |\n|-----------|-------------|--------------|\n| `for`    | Known repeats | Iterates over items |\n| `while`  | Unknown repeats | Runs till condition breaks |",
        "slide_type": "table"
      },
      {
        "title": "Want to Loop Like a Pro? 🤓",
        "content": "Next video: How to break out of loops early like a ninja using `break` and `continue`! 🥷🔥",
        "slide_type": "text"
      }
    ]
  }
]






script text:

[
  {
    "topic": "if/elif/else",
    "slides": [
      {
        "title": "Stop Writing Messy if/else Statements in Python!",
        "script": "Have you ever looked at your if/else code and thought, \"This is confusing!\"? [pause] Don’t worry, you’re not alone. Messy blocks make your program hard to read and full of bugs. [pause] But there’s a simple trick to keep things clean and easy. Let’s dive into making your code neat and smooth to follow."
      },
      {
        "title": "Why Use `elif`?",
        "script": "Why use elif instead of many ifs? [pause] Imagine checking your exam score to decide your grade. Instead of separate checks, elif links conditions in one flow. This way, Python stops checking as soon as one condition matches. [pause] It saves time and makes your code clearer. Like choosing the right path in a road junction, elif guides your program to the right place."
      },
      {
        "title": "Common Mistake 🚫",
        "script": "A common mistake is using many ifs for related checks. [pause] This causes your program to run multiple blocks, even if one condition is already true. That slows things down and makes debugging tricky. Think of it like asking your neighbor the same question again and again—it wastes time. Using elif prevents this by making sure only one block runs."
      },
      {
        "title": "Quick Tips ✨",
        "script": "Here’s a quick way to remember if/elif/else. Start with if for your first check. Use elif for all the middle steps. Finally, add else as a backup plan if none match. [pause] It’s like a flowchart where your code decides the best path to take. This keeps your program neat and easy to follow."
      },
      {
        "title": "Want Cleaner Code? 👀",
        "script": "Now that you know how to clean your if/else statements, what’s next? [pause] Soon, we’ll explore ternary operators in Python. These one-liners make your code even neater and can impress any interviewer. For now, keep practicing these basics—they’re the foundation of great coding!"
      }
    ]
  },
  {
    "topic": "for vs while loops",
    "slides": [
      {
        "title": "Python for vs while Loops: The 60-Second Rule",
        "script": "Loops help you repeat tasks in your program. But when should you pick for or while? [pause] Here’s a quick rule to remember. This will help you decide fast, so you don’t get stuck choosing the wrong loop type. Ready? Let’s start!"
      },
      {
        "title": "Use `for` When...",
        "script": "Use for loops when you know exactly how many times to repeat. For example, if you want to print numbers 0 to 4, for loop does it smoothly. [pause] It works perfectly with lists or ranges, like counting steps in a recipe. Think of it as following a fixed set of instructions."
      },
      {
        "title": "Use `while` When...",
        "script": "While loops are best when you don’t know how many times to repeat. Instead, you keep going until a condition stops being true. For example, counting until a number reaches five. [pause] It’s like waiting for rain to stop before going outside—you don’t know exactly when, but you keep checking."
      },
      {
        "title": "Quick Comparison Table 📊",
        "script": "Let’s compare for and while loops quickly. For loops repeat a known number of times by going through each item. While loops run until a condition is false, so they can run unknown times. [pause] Remember this: for is like walking fixed steps, while is like walking until you find a shop."
      },
      {
        "title": "Want to Loop Like a Pro? 🤓",
        "script": "You now understand when to use for and while loops. [pause] Next, we’ll learn how to control loops better using break and continue. These help you stop or skip steps in your loops like a coding ninja. For now, practice these ideas to build strong basics."
      }
    ]
  }
]



I will provide the script text along with ppt json, can you please same thing you can add script to create the md file using the button and one textarea add for the inputing the json and once the click to process it will generate the script file with the inside that topic folder and give me the final app code

